"""Generate VKS on VCF 9.1 presentation (.pptx) from README content and screenshots."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import os

SCREENSHOTS = os.path.join(os.path.dirname(__file__), 'screenshots')
OUT = os.path.join(os.path.dirname(__file__), 'VKS-on-VCF91.pptx')

# VMware brand colours
DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT    = RGBColor(0x00, 0xB0, 0xEA)   # VMware blue
GREEN     = RGBColor(0x4C, 0xAF, 0x50)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank


# ─── helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)   # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             wrap=True, v_anchor=None):
    from pptx.enum.text import MSO_ANCHOR
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_image_safe(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        # placeholder box
        add_rect(slide, left, top, width, height, RGBColor(0xCC, 0xCC, 0xCC))
        add_text(slide, os.path.basename(path), left, top, width, height,
                 font_size=10, color=DARK_TEXT, align=PP_ALIGN.CENTER)


def dark_slide(title, subtitle=None):
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, W, H, DARK_BG)
    # accent bar left
    add_rect(slide, 0, 0, Inches(0.12), H, ACCENT)
    add_text(slide, title,
             Inches(0.3), Inches(2.5), Inches(12.7), Inches(1.8),
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.3), Inches(4.4), Inches(12.7), Inches(1.2),
                 font_size=22, bold=False, color=ACCENT, align=PP_ALIGN.LEFT)
    return slide


def content_slide(title):
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, W, H, LIGHT_BG)
    # header bar
    add_rect(slide, 0, 0, W, Inches(1.0), DARK_BG)
    add_rect(slide, 0, 0, Inches(0.12), Inches(1.0), ACCENT)
    add_text(slide, title,
             Inches(0.25), Inches(0.12), Inches(12.7), Inches(0.76),
             font_size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    return slide


def bullet_block(slide, items, left, top, width, height, font_size=16, indent=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        prefix = '    • ' if indent else '• '
        run.text = prefix + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = DARK_TEXT


def two_col_table(slide, headers, rows, left, top, width, height, font_size=14):
    """Simple two-column comparison rendered as coloured boxes."""
    n_cols = len(headers)
    col_w = width // n_cols
    row_h = height // (len(rows) + 1)

    # header
    for ci, h in enumerate(headers):
        bg = DARK_BG if ci == 0 else ACCENT
        add_rect(slide, left + ci * col_w, top, col_w, row_h, bg)
        add_text(slide, h,
                 left + ci * col_w + Inches(0.05), top, col_w, row_h,
                 font_size=font_size, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    for ri, row in enumerate(rows):
        y = top + (ri + 1) * row_h
        for ci, cell in enumerate(row):
            bg = RGBColor(0xE8, 0xF4, 0xFB) if ci % 2 == 0 else WHITE
            add_rect(slide, left + ci * col_w, y, col_w, row_h, bg,
                     line_rgb=RGBColor(0xCC, 0xCC, 0xCC))
            add_text(slide, str(cell),
                     left + ci * col_w + Inches(0.05), y + Pt(2), col_w, row_h,
                     font_size=font_size - 1, bold=False, color=DARK_TEXT)


# ─── SLIDE 1 — Cover ─────────────────────────────────────────────────────────
slide = dark_slide(
    'VKS on VCF 9.1',
    'vSphere Kubernetes Service 端到端自動化\nVMware Cloud Foundation 9.1 (rtolab)'
)
add_text(slide, '✅  2026-06-08 實機跑通  •  Path A (DTGW+VNA)  •  v1.34.2',
         Inches(0.3), Inches(5.6), Inches(12.7), Inches(0.6),
         font_size=14, bold=False, color=GREEN, align=PP_ALIGN.LEFT)

# ─── SLIDE 2 — 目標與四種方法 ─────────────────────────────────────────────────
slide = content_slide('目標：四種方法起 VKS')
items = [
    'UI（畫面）— vCenter Activate Supervisor wizard + Namespace + VKS cluster，7 步，逐步截圖',
    'API（raw REST / curl）— vCenter /api、NSX /policy/api/v1、Supervisor kube API',
    'Python — requests 打 REST，step1~4 腳本，lab.py 共用設定，支援環境變數覆寫',
    'PowerShell — PowerCLI + REST，Step0~4.ps1，-DryRun 支援，生產可直接跑',
]
bullet_block(slide, items, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5), font_size=17)

# ─── SLIDE 3 — 兩條路線 ──────────────────────────────────────────────────────
slide = content_slide('兩條網路路線')
headers = ['項目', 'Path A — DTGW + VNA', 'Path B — Edge + Centralized TGW']
rows = [
    ['Transit GW 跑在', 'ESXi host（Distributed）', 'NSX Edge VM（Centralized）'],
    ['需要 NSX Edge cluster', '❌ 不需要', '✅ 需要（2× MEDIUM edge VM）'],
    ['Stateful 服務（SNAT/LB）', 'VNA cluster（9.1 新功能）', 'Edge node'],
    ['Lab 現狀', '預設 DTGW，只缺 VNA → ✅ 部署', 'Edge cluster 未部 → 待辦'],
    ['部署成本', '部 1 個 VNA cluster', '部 edge cluster + 改 TGW span'],
    ['本次實測', '✅ 端到端成功', '—（lab edge 未部，未測）'],
]
two_col_table(slide, headers, rows,
              Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5), font_size=13)

# ─── SLIDE 4 — Lab NSX 實況 ──────────────────────────────────────────────────
slide = content_slide('Lab NSX 實況（2026-06-08）')
left_items = [
    'NSX 版本：9.1.0.0.25318225',
    'vCenter：192.168.114.11 (kosten-vcf91-vc.rtolab.local)',
    'NSX Manager VIP：192.168.114.13',
    'Default TGW：存在，ClusterBasedSpan = DTGW，transit subnet 100.64.0.0/21',
    'Default VPC Connectivity Profile：存在，指向 default DTGW',
    'NSX Edge cluster：0 個（inventory 規劃 vcf-m02-edge-cl01 未部）',
    'VNA cluster：0 個（feature 在，未部署）→ Path A 需先部',
]
bullet_block(slide, left_items, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5), font_size=16)

# ─── SLIDE 5 — IP 規劃 ───────────────────────────────────────────────────────
slide = content_slide('IP 規劃（rtolab 192.168.114.0/24）')
headers = ['用途', 'IP / CIDR', '備註']
rows = [
    ['Supervisor CP VMs (5 IPs)', '192.168.114.101–105', '兩路線通用'],
    ['Supervisor mgmt gateway', '192.168.114.254', '兩路線通用'],
    ['NSX External IP Block (LB/SNAT)', '192.168.114.128/26', '/26 邊界對齊！'],
    ['Private TGW IP Block', '172.30.0.0/16', '/16 硬規定；避開 VCFA 172.27/16'],
    ['Supervisor Service CIDR', '172.29.0.0/16', 'K8s ClusterIP（Supervisor 用）'],
    ['VPC Default Private CIDR', '172.28.0.0/16', 'namespace 子網來源（不可與 TGW block 重疊）'],
    ['VKS cluster Pod CIDR', '100.96.0.0/11', '避開 192.168.114.0/24 管理網段！'],
    ['VKS cluster Service CIDR', '10.96.0.0/12', 'per-cluster'],
    ['K8s API Server (Supervisor)', '192.168.114.132', '從 External Block 分配'],
    ['VKS CP 外部 IP', '192.168.114.135', '從 External Block 分配'],
]
two_col_table(slide, headers, rows,
              Inches(0.3), Inches(1.1), Inches(12.7), Inches(6.0), font_size=13)

# ─── SLIDE 6 — 執行流程 ──────────────────────────────────────────────────────
slide = content_slide('執行流程')
steps = [
    'Step 0  common/Step0-Check-Prereqs.ps1  — 前置條件確認（兩路線都先跑）',
    'Step 1  path-a-dtgw/Step1-Setup-DTGW.ps1  — NSX IP blocks + VPC Connectivity Profile + VNA cluster',
    'Step 2  common/Step2-Enable-Supervisor.ps1  — 啟用 Supervisor（VPC mode）',
    'Step 3  common/Step3-New-Namespace.ps1  — 建 namespace（storage policy + access）',
    'Step 4  common/Step4-New-VksCluster.ps1  — 建 VKS guest cluster（Cluster CR via kube API）',
    '          → kubeconfig 輸出，給 automation 直接用',
]
bullet_block(slide, steps, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5), font_size=16)

# small note
add_text(slide, '所有 Step*.ps1 支援 -DryRun（印 payload 不送出）',
         Inches(0.3), Inches(6.7), Inches(12.7), Inches(0.5),
         font_size=13, bold=False, color=RGBColor(0x55, 0x55, 0x55))

# ─── SLIDE 7 — NSX IP Address Blocks ─────────────────────────────────────────
slide = content_slide('NSX 前置設定：IP Address Blocks')
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '40-nsx-ip-blocks.jpg'),
               Inches(0.3), Inches(1.1), Inches(8.5), Inches(5.8))
add_rect(slide, Inches(9.1), Inches(1.1), Inches(4.0), Inches(5.8), DARK_BG)
notes = [
    'External Block（SNAT/LB 來源）',
    '  vcf-m02-vks-ext-ipblock',
    '  192.168.114.128/26',
    '  /26 邊界對齊！(.128,.192,.0)',
    '',
    'Private TGW Block（VKS 必要）',
    '  vcf-m02-vks-priv-tgw',
    '  172.30.0.0/16',
    '  /16 硬規定',
    '',
    '(系統預設 default--kube-s...',
    '  172.28.0.0/16 勿動)',
]
y = Inches(1.3)
for n in notes:
    add_text(slide, n, Inches(9.2), y, Inches(3.7), Inches(0.32),
             font_size=12, bold=False, color=WHITE)
    y += Inches(0.37)

# ─── SLIDE 8 — NSX VPC Connectivity Profile ───────────────────────────────────
slide = content_slide('NSX 前置設定：VPC Connectivity Profile')
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '41-nsx-vpc-profiles.jpg'),
               Inches(0.3), Inches(1.1), Inches(8.5), Inches(5.8))
add_rect(slide, Inches(9.1), Inches(1.1), Inches(4.0), Inches(5.8), DARK_BG)
notes2 = [
    'vcf-m02-vks-vpc-profile',
    '  Transit Gateway:',
    '    Default Transit Gateway',
    '    (DTGW / distributed)',
    '',
    '  External IP Blocks:',
    '    vcf-m02-vks-ext-ipblock',
    '',
    '  Private TGW IP Blocks:',
    '    vcf-m02-vks-priv-tgw',
    '',
    '  Status: Success ✅',
]
y = Inches(1.3)
for n in notes2:
    add_text(slide, n, Inches(9.2), y, Inches(3.7), Inches(0.32),
             font_size=12, bold=False, color=WHITE)
    y += Inches(0.37)

# ─── SLIDE 8b — Wizard Section Header ────────────────────────────────────────
slide = dark_slide(
    'Activate Supervisor Wizard\n實機 UI Walkthrough',
    'Steps 1–7 全部截圖（2026-06-08 Step1-3；2026-06-09 Step4-7）'
)

# ─── SLIDE 8c — Wizard 入口 + Step 1 ─────────────────────────────────────────
slide = content_slide('Wizard 入口 & Step 1 — vCenter + 網路模式')
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '00-intro.jpg'),
               Inches(0.3), Inches(1.1), Inches(6.2), Inches(5.6))
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '01-step1-vcenter-network.jpg'),
               Inches(6.7), Inches(1.1), Inches(6.3), Inches(5.6))
add_text(slide, '⚠️ networking stack 只有「VCF Networking with VPC」和「VDS」兩個，9.1 沒有獨立 NSX-classic 選項',
         Inches(0.3), Inches(6.8), Inches(12.7), Inches(0.5),
         font_size=12, bold=False, color=DARK_TEXT)

# ─── SLIDE 8d — Wizard Step 2 ─────────────────────────────────────────────────
slide = content_slide('Wizard Step 2 — Supervisor 位置（Cluster Deployment）')
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '06-step2-name-cluster-selected.jpg'),
               Inches(0.3), Inches(1.1), Inches(8.5), Inches(5.8))
add_rect(slide, Inches(9.1), Inches(1.1), Inches(4.0), Inches(5.8), DARK_BG)
s2_notes = [
    'CLUSTER DEPLOYMENT tab',
    '',
    'Supervisor name:',
    '  vcf-m02-supervisor',
    'Cluster: vcf-m02-cl01',
    '  COMPATIBLE ✅',
    '  4 hosts',
    '  CPU: 137.04 GHz',
    '  Memory: 171.51 GB',
    '',
    '⚠️ 先點 datacenter 節點',
    '  vcf-m02-dc，COMPATIBLE',
    '  tab 才會列出 cluster',
]
y = Inches(1.3)
for n in s2_notes:
    add_text(slide, n, Inches(9.2), y, Inches(3.7), Inches(0.33),
             font_size=12, bold=False, color=WHITE)
    y += Inches(0.37)

# ─── SLIDE 8e — Wizard Step 3 ─────────────────────────────────────────────────
slide = content_slide('Wizard Step 3 — Storage Policy')
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '07-step3-storage.jpg'),
               Inches(0.3), Inches(1.1), Inches(8.5), Inches(5.8))
add_rect(slide, Inches(9.1), Inches(1.1), Inches(4.0), Inches(5.8), DARK_BG)
s3_notes = [
    '三個 policy 全選：',
    '  Management Storage',
    '  Policy - Single Node',
    '  (FTT=0，lab 用)',
    '',
    '• Control Plane Storage Policy',
    '• Ephemeral Disks Storage Policy',
    '• Image Cache Storage Policy',
    '',
    '⚠️ SPBM dropdown 在 nested',
    '  vCenter 載入 ~2–10 分鐘',
    '  → renderer 凍結',
    '  → Step 4-7 於隔日補截圖',
]
y = Inches(1.3)
for n in s3_notes:
    add_text(slide, n, Inches(9.2), y, Inches(3.7), Inches(0.35),
             font_size=12, bold=False, color=WHITE)
    y += Inches(0.37)

# ─── SLIDE 8f — Wizard Step 4 Management Network ─────────────────────────────
slide = content_slide('Wizard Step 4 — Management Network')
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '08-step4-mgmt-network.jpg'),
               Inches(0.3), Inches(1.1), Inches(8.5), Inches(5.8))
add_rect(slide, Inches(9.1), Inches(1.1), Inches(4.0), Inches(5.8), DARK_BG)
s4_notes = [
    'IP Assignment Mode: Static',
    'domain-c9 Network:',
    '  vcf-m02-cl01-vds01-pg-mgmt',
    '',
    'IP Addresses:',
    '  192.168.114.101 – 105',
    'Subnet Mask: 255.255.255.0',
    'Gateway: 192.168.114.254',
    'DNS Server(s): 192.168.114.200',
    'DNS Search: rtolab.local',
    'NTP Server(s): 192.168.114.200',
]
y = Inches(1.3)
for n in s4_notes:
    add_text(slide, n, Inches(9.2), y, Inches(3.7), Inches(0.38),
             font_size=12, bold=False, color=WHITE)
    y += Inches(0.38)

# ─── SLIDE 8g — Wizard Step 5 Workload Network ───────────────────────────────
slide = content_slide('Wizard Step 5 — Workload Network（VPC mode）')
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '09-step5-workload-network.jpg'),
               Inches(0.3), Inches(1.1), Inches(8.5), Inches(5.8))
add_rect(slide, Inches(9.1), Inches(1.1), Inches(4.0), Inches(5.8), DARK_BG)
s5_notes = [
    'NSX Project: Default',
    'VPC Connectivity Profile:',
    '  vcf-m02-vks-vpc-profile ✅',
    '',
    'External IP Blocks:',
    '  vcf-m02-vks-ext-ipblock',
    '  192.168.114.128/26',
    '',
    'Private TGW IP Blocks:',
    '  vcf-m02-vks-priv-tgw',
    '  172.30.0.0/16',
    '',
    'Service CIDR: 172.29.0.0/16',
    'DNS/NTP: 192.168.114.200',
]
y = Inches(1.3)
for n in s5_notes:
    add_text(slide, n, Inches(9.2), y, Inches(3.7), Inches(0.35),
             font_size=12, bold=False, color=WHITE)
    y += Inches(0.35)

# ─── SLIDE 8h — Wizard Step 6 Advanced Settings ──────────────────────────────
slide = content_slide('Wizard Step 6 — Advanced Settings')
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '10-step6-advanced.jpg'),
               Inches(0.3), Inches(1.1), Inches(8.5), Inches(5.8))
add_rect(slide, Inches(9.1), Inches(1.1), Inches(4.0), Inches(5.8), DARK_BG)
s6_notes = [
    'Supervisor Control Plane Size:',
    '  Small',
    '  (CPUs: 4 / Memory: 16 GB',
    '   Storage: 48 GB)',
    '',
    'API Server DNS Name(s):',
    '  Optional（lab 留空）',
    '',
    'Export configuration:',
    '  Unchecked',
    '',
    '→ 無 Content Library 選項',
    '  （Wizard 另有 Prerequisites',
    '  步驟指定 Library）',
]
y = Inches(1.3)
for n in s6_notes:
    add_text(slide, n, Inches(9.2), y, Inches(3.7), Inches(0.36),
             font_size=12, bold=False, color=WHITE)
    y += Inches(0.36)

# ─── SLIDE 8i — Wizard Step 7 Ready to Complete ──────────────────────────────
slide = content_slide('Wizard Step 7 — Ready to Complete')
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '11-step7-ready.jpg'),
               Inches(0.3), Inches(1.1), Inches(8.5), Inches(5.8))
add_rect(slide, Inches(9.1), Inches(1.1), Inches(4.0), Inches(5.8), DARK_BG)
s7_notes = [
    'vCenter Server:',
    '  kosten-vcf91-vc.rtolab.local',
    'Network: VCF Networking with VPC',
    '',
    'Supervisor Name:',
    '  vcf-m02-supervisor',
    'vSphere Zone: domain-c9',
    'CP HA: Disabled',
    '',
    'Mgmt IP: .101–.105',
    'GW: .254 / DNS/NTP: .200',
    '',
    '⚠️ 未按 FINISH',
    '  （截圖目的，已另用 API 部署）',
]
y = Inches(1.3)
for n in s7_notes:
    add_text(slide, n, Inches(9.2), y, Inches(3.7), Inches(0.36),
             font_size=12, bold=False, color=WHITE)
    y += Inches(0.36)

# ─── SLIDE 9 — Supervisor Configure Network (Management) ──────────────────────
slide = content_slide('Supervisor Configure — 管理網路（Management Network）')
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '43-vc-sup-configure-network-mgmt.jpg'),
               Inches(0.3), Inches(1.1), Inches(8.5), Inches(5.8))
add_rect(slide, Inches(9.1), Inches(1.1), Inches(4.0), Inches(5.8), DARK_BG)
mgmt_notes = [
    'IP Assignment Mode: Static',
    'Network: vcf-m02-cl01-vds01-pg-mgmt',
    'Starting IP: 192.168.114.101',
    '  (5 IPs → .101 ~ .105)',
    'Subnet Mask: 255.255.255.0',
    'Gateway: 192.168.114.254',
    'DNS: 192.168.114.200',
    'DNS Domain: rtolab.local',
    'NTP: 192.168.114.200',
]
y = Inches(1.3)
for n in mgmt_notes:
    add_text(slide, n, Inches(9.2), y, Inches(3.7), Inches(0.35),
             font_size=12, bold=False, color=WHITE)
    y += Inches(0.38)

# ─── SLIDE 10 — Supervisor Configure Network (Workload) ───────────────────────
slide = content_slide('Supervisor Configure — 工作負載網路（Workload Networks）')
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '43b-vc-sup-configure-network-workload.jpg'),
               Inches(0.3), Inches(1.1), Inches(8.5), Inches(5.8))
add_rect(slide, Inches(9.1), Inches(1.1), Inches(4.0), Inches(5.8), DARK_BG)
wl_notes = [
    'NSX Project: Default',
    'VPC Connectivity Profile:',
    '  vcf-m02-vks-vpc-profile',
    '',
    'External IP Blocks:',
    '  vcf-m02-vks-ext-ipblock',
    '  192.168.114.128/26',
    '  Usage: 12.5%',
    '',
    'Private TGW IP Blocks:',
    '  vcf-m02-vks-priv-tgw',
    '  172.30.0.0/16',
    '  Usage: 0.02%',
]
y = Inches(1.3)
for n in wl_notes:
    add_text(slide, n, Inches(9.2), y, Inches(3.7), Inches(0.32),
             font_size=12, bold=False, color=WHITE)
    y += Inches(0.37)

# ─── SLIDE 11 — Supervisor Configure Storage ──────────────────────────────────
slide = content_slide('Supervisor Configure — 儲存設定（Storage）')
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '44-vc-sup-configure-storage.jpg'),
               Inches(0.3), Inches(1.1), Inches(8.5), Inches(5.8))
add_rect(slide, Inches(9.1), Inches(1.1), Inches(4.0), Inches(5.8), DARK_BG)
stor_notes = [
    'Control Plane Nodes:',
    '  Management Storage',
    '  Policy - Single Node',
    '  (FTT=0, nested OK)',
    '',
    'Ephemeral Disks:',
    '  同上',
    '',
    'Image Cache:',
    '  同上',
    '',
    '⚠️ 生產環境建議用 FTT=1+',
    '  (此為 lab single-node)',
]
y = Inches(1.3)
for n in stor_notes:
    add_text(slide, n, Inches(9.2), y, Inches(3.7), Inches(0.32),
             font_size=12, bold=False, color=WHITE)
    y += Inches(0.37)

# ─── SLIDE 12 — Supervisor Running 截圖 ───────────────────────────────────────
slide = content_slide('Supervisor 運行狀態')
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '30-supervisor-running.jpg'),
               Inches(0.3), Inches(1.1), Inches(12.7), Inches(6.1))

# ─── SLIDE 8 — Content Libraries ─────────────────────────────────────────────
slide = content_slide('必備：兩個 Content Library')
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '34-content-library-list.jpg'),
               Inches(0.3), Inches(1.1), Inches(8.0), Inches(5.5))
# caption box
add_rect(slide, Inches(8.5), Inches(1.1), Inches(4.5), Inches(5.5), DARK_BG)
bullets = [
    'tkg-content-library',
    '  → Supervisor image',
    '  → 訂閱 .../supervisor/v1/latest/lib.json',
    '',
    'tkg-tkr-library',
    '  → TKG node image (Photon OS)',
    '  → 訂閱 .../v2/latest/lib.json',
    '  → 7.69 GB / 123 items',
    '',
    '⚠️ 少了 tkr-library',
    '  tkr-resolver 報',
    '  "Could not resolve KR/OSImage"',
    '  → cluster 建不起來',
]
y = Inches(1.3)
for b in bullets:
    add_text(slide, b, Inches(8.6), y, Inches(4.2), Inches(0.35),
             font_size=12, bold=False, color=WHITE)
    y += Inches(0.33)

# ─── SLIDE 9 — Namespace 截圖 ─────────────────────────────────────────────────
slide = content_slide('Namespace vks-automation')
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '33-namespace-vks-automation.jpg'),
               Inches(0.3), Inches(1.1), Inches(6.2), Inches(5.5))
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '32-namespaces-list.jpg'),
               Inches(6.7), Inches(1.1), Inches(6.3), Inches(5.5))

# ─── SLIDE 10 — VKS Cluster Running ─────────────────────────────────────────
slide = content_slide('VKS Cluster 建立完成')
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '37-vks-cluster-running.jpg'),
               Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.5))
add_text(slide, 'vks-auto-01  •  Available = True  •  v1.34.2+vmware.2  •  CP 192.168.114.135',
         Inches(0.3), Inches(6.75), Inches(12.7), Inches(0.5),
         font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ─── SLIDE 11 — VKS CP + Worker VMs ─────────────────────────────────────────
slide = content_slide('CP + Worker VM 狀態')
add_image_safe(slide,
               os.path.join(SCREENSHOTS, '36-vks-cluster-vms.jpg'),
               Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.5))
add_text(slide, 'CP: vks-auto-01-7f6ms-fmgz9  •  Worker: vks-auto-01-node-pool-1-5cjwv-...',
         Inches(0.3), Inches(6.75), Inches(12.7), Inches(0.5),
         font_size=14, bold=False, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# ─── SLIDE 12 — 踩坑 / 必修正 ─────────────────────────────────────────────────
slide = content_slide('踩坑必修正（否則建不起來）')
items_left = [
    'Pod CIDR 不能用預設 192.168.0.0/16',
    '  → 與管理網段 192.168.114.0/24 重疊',
    '  → kubeadm init 後 Antrea overlay route',
    '    劫持 DNS / Registry 流量 → CP never Ready',
    '  ✅ 改用 100.96.0.0/11',
]
items_right = [
    'MachineHealthCheck timeout 太短',
    '  → 預設 nodeStartupTimeout = 3600s',
    '  → nested-on-nested 啟動 >60min',
    '  → MHC 提早 remediate → 砍 CP → 死循環',
    '  ✅ topology.controlPlane.healthCheck',
    '     .checks.nodeStartupTimeoutSeconds: 14400',
]
for item in items_left:
    pass
add_rect(slide, Inches(0.3), Inches(1.2), Inches(6.2), Inches(5.5), DARK_BG)
add_rect(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.5), DARK_BG)

y = Inches(1.4)
for b in items_left:
    add_text(slide, b, Inches(0.45), y, Inches(6.0), Inches(0.4),
             font_size=13, bold=False, color=WHITE)
    y += Inches(0.38)

y = Inches(1.4)
for b in items_right:
    add_text(slide, b, Inches(6.95), y, Inches(6.0), Inches(0.4),
             font_size=13, bold=False, color=WHITE)
    y += Inches(0.38)

# labels
add_text(slide, '① Pod CIDR 衝突',
         Inches(0.3), Inches(1.15), Inches(3.0), Inches(0.35),
         font_size=14, bold=True, color=ACCENT)
add_text(slide, '② MachineHealthCheck Timeout',
         Inches(6.8), Inches(1.15), Inches(5.0), Inches(0.35),
         font_size=14, bold=True, color=ACCENT)

# ─── SLIDE 13 — NSX SNAT 驗證 ────────────────────────────────────────────────
slide = content_slide('NSX SNAT 驗證通過')
facts = [
    '實測：cluster 建立過程 NSX SNAT session 有效',
    'SNAT bytes out：99,960 packets / 144 MB（拉 container image）',
    'SNAT source：192.168.114.128/26 external block',
    'VNA cluster 負責 SNAT/LB（Path A DTGW 路線）',
    '',
    'kubectl get nodes 輸出：',
    '  vks-auto-01-7f6ms-fmgz9       Ready   control-plane   v1.34.2+vmware.2',
    '  vks-auto-01-node-pool-1-...   Ready   <none>          v1.34.2+vmware.2',
    '',
    'Container runtime：containerd 2.1.5-fips',
    'OS：Photon OS (Linux)',
]
add_rect(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5), DARK_BG)
y = Inches(1.4)
for f in facts:
    col = ACCENT if f.startswith('kubectl') or f.startswith('  vks') else WHITE
    add_text(slide, f, Inches(0.5), y, Inches(12.3), Inches(0.4),
             font_size=14, bold=False, color=col)
    y += Inches(0.38)

# ─── SLIDE 14 — Python method ────────────────────────────────────────────────
slide = content_slide('Python 方法（開發用）')
code_lines = [
    '# 安裝',
    'pip install -r requirements.txt   # requests, PyYAML, kubernetes',
    '',
    '# 設定（或 export VC=... VCPASS=... NSXPASS=...）',
    '# python/lab.py  — 共用連線設定',
    '',
    '# Step 1 — NSX IP blocks + VPC profile',
    'python step1_setup_dtgw.py --dry-run',
    'python step1_setup_dtgw.py',
    '',
    '# Step 2 — 啟用 Supervisor',
    'python step2_enable_supervisor.py',
    '',
    '# Step 3 — 建 namespace',
    'python step3_new_namespace.py',
    '',
    '# Step 4 — 建 VKS cluster (kubectl 或 kubernetes client)',
    'kubectl-vsphere login --server=192.168.114.132 -u administrator@vsphere.local --insecure-skip-tls-verify',
    'python step4_new_vks_cluster.py',
    'python step4_new_vks_cluster.py --client   # 純 python-kubernetes',
]
add_rect(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.8), DARK_BG)
y = Inches(1.35)
for line in code_lines:
    col = ACCENT if line.startswith('#') else WHITE
    add_text(slide, line, Inches(0.5), y, Inches(12.3), Inches(0.27),
             font_size=11.5, bold=False, color=col)
    y += Inches(0.265)

# ─── SLIDE 15 — PowerShell / API 方法 ────────────────────────────────────────
slide = content_slide('PowerShell & API 方法')
add_rect(slide, Inches(0.3), Inches(1.2), Inches(6.0), Inches(5.5), DARK_BG)
add_rect(slide, Inches(6.6), Inches(1.2), Inches(6.4), Inches(5.5), DARK_BG)

ps_lines = [
    '# PowerShell（common/ + path-a-dtgw/）',
    '',
    'pwsh Step0-Check-Prereqs.ps1',
    'pwsh path-a-dtgw/Step1-Setup-DTGW.ps1 \\',
    '     -DryRun    # 先 dry-run',
    'pwsh path-a-dtgw/Step1-Setup-DTGW.ps1',
    'pwsh Step2-Enable-Supervisor.ps1',
    'pwsh Step3-New-Namespace.ps1',
    'pwsh Step4-New-VksCluster.ps1',
    '',
    '# PowerCLI cmdlet 路線',
    'Enable-WMCluster ...',
    'New-WMNamespace ...',
]
y = Inches(1.4)
for line in ps_lines:
    col = ACCENT if line.startswith('#') else WHITE
    add_text(slide, line, Inches(0.5), y, Inches(5.7), Inches(0.33),
             font_size=11, bold=False, color=col)
    y += Inches(0.32)

api_lines = [
    '# API（curl / 語言無關）',
    '',
    '# 認證',
    'SID=$(curl -u admin:pass -X POST',
    '  https://vc/api/session)',
    '',
    '# NSX IP blocks',
    'curl -X PATCH https://nsx/policy/api/v1',
    '  /infra/ip-blocks/vcf-m02-vks-ext',
    '',
    '# 啟用 Supervisor',
    'curl -X POST https://vc/api/vcenter',
    '  /namespace-management/supervisors',
    '',
    '# VKS cluster (Cluster CR)',
    'kubectl apply -f common/vks-cluster.yaml',
]
y = Inches(1.4)
for line in api_lines:
    col = ACCENT if line.startswith('#') else WHITE
    add_text(slide, line, Inches(6.75), y, Inches(6.1), Inches(0.33),
             font_size=11, bold=False, color=col)
    y += Inches(0.32)

# ─── SLIDE 16 — Closing / Links ──────────────────────────────────────────────
slide = dark_slide('完成', 'VKS on VCF 9.1 — 端到端自動化')
links = [
    'GitHub Repo：github.com/kostenyang/vcf9.1vks',
    'README.md  →  四種方法、IP 規劃、執行流程',
    'research/05-test-execution.md  →  踩坑與修正完整紀錄',
    'python/  →  Step1~4 Python 腳本（requests + kubernetes client）',
    'common/vks-cluster.yaml  →  ClusterClass CR（含 Pod CIDR + MHC 修正）',
    'screenshots/  →  38 張實機截圖（Wizard Step1-7 + NSX + Supervisor Configure 前置設定）',
]
y = Inches(4.5)
for link in links:
    add_text(slide, '→  ' + link, Inches(0.3), y, Inches(12.7), Inches(0.45),
             font_size=16, bold=False, color=ACCENT)
    y += Inches(0.43)

# ─── Save ────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f'Saved: {OUT}')
print(f'Slides: {len(prs.slides)}')
